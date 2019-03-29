# -*- coding: utf-8 -*-

from gensim.models import Doc2Vec
from title_master import title_master
import numpy as np
import pickle
from datetime import datetime
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

criteria_score = 0.4
criteria_num = 5

title_master_class = title_master()
title_dict = title_master_class.get_dict()
title_dict_keys = list(title_dict.keys())

archive_frequency_word = pickle.load(open('archive_frequency_word_by10.pickle', 'rb'))
model = Doc2Vec.load('doc2vec_merge.model')

# 類似度行列
similarity_matrix = np.zeros((92, 92))

for i in range(92):
    for k in range(92):
        similarity_matrix[i, k] = round(model.docvecs.similarity(i, k), 3)

today = datetime.today().strftime("%Y%m%d")
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "アーカイブ論文類似度"
subtitle.text = str(today) + "堀智之"
for i in range(92):
    index_list = np.argsort(-similarity_matrix[i, :])
    
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    t = ''
    for w in archive_frequency_word[title_dict_keys[i]]:
        t = t + '「' + w + '」'
    title_shape.text = "論文名：" + title_dict[title_dict_keys[i]] + '\n' + '頻出単語' + t
    title_shape.text_frame.paragraphs[0].font.size = Pt(22)  # font size
    title_shape.text_frame.paragraphs[0].font.bold = True  # font bold
    
    tf = body_shape.text_frame
    tf.text = "類似論文と頻出単語"
    tf.paragraphs[0].font.size = Pt(18)  # font size
    tf.paragraphs[0].font.bold = True  # font bold
    for k in range(len(index_list)):
        index = index_list[k]
        if k == 0:
            continue
        if k == criteria_num + 1:
            break
        p = tf.add_paragraph()
        p.text = title_dict[title_dict_keys[index]] + '(類似度' + str(round(similarity_matrix[i, index], 3)) + ')'
        p.font.size = Pt(16)  # font size
        p.font.bold = True  # font bold
        if similarity_matrix[i, index] > criteria_score:
            p.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
        p.level = 1
        p = tf.add_paragraph()
        t = ''
        for w in archive_frequency_word[title_dict_keys[index]]:
            t = t + '「' + w + '」'
        p.text = '頻出単語' + t
        p.font.size = Pt(14)  # font size
        p.font.bold = True  # font bold
        p.level = 2
    
prs.save("archive%s.pptx" % today)