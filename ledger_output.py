import openpyxl
from openpyxl.styles import PatternFill, Font, Alignment, Border, Side
from openpyxl.drawing.image import Image
import sys
import importlib
importlib.reload(sys)

from pdfminer.pdfparser import PDFParser,PDFDocument
from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from pdfminer.converter import PDFPageAggregator
from pdfminer.layout import *
from pdfminer.pdfinterp import PDFTextExtractionNotAllowed
from PyPDF2 import PdfFileReader
import fitz
import os
import xlsxwriter
import re
import extract_img_from_pdf as eifp
import cv2
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
import text_analysis as ta

DATA_PATH = './data/'
IMG_PATH = './IMG/'
TEMP_PATH = './temp/'

dic_region_type = {1: '部', 
                   2: '章', 
                   3: '節', 
                   4: '小節', 
                   5: '段落', 
                   6: '図タイトル'}

class TextRegion:
    """Text region structure

    Arguments:
        region_type: Integer, type of text region (1: Part, 2: Chapter, 3: Section, 4: Sub-section, 5: Paragraph, 6: Title, 7: Page number)
        x0: Float, x-coordinate in top-left
        y0: Float, y-coordinate in top-left
        x1: Float, x-coordinate in bottom-right
        y1: Float, y-coordinate in bottom-right
        line_height: Float, average height of text
        text: String, text
        page_number
        part_number
        chapter_number
        section_number
        sub_section_number
        paragraph_number
        title_number
        img_path
        is_img_text
        text_file_path
        excel_row

    """
    

    def __init__(self, region_type, x0, y0, x1, y1, line_height, text):
        self.region_type = region_type
        self.x0 = x0
        self.y0 = y0
        self.x1 = x1
        self.y1 = y1
        self.line_height = line_height
        self.text = text
        self.page_number = -1
        self.part_number = -1
        self.chapter_number = -1
        self.section_number = -1
        self.sub_section_number = -1
        self.paragraph_number = -1
        self.title_number = ''
        self.img_path = ''
        self.is_img_text = False
        self.text_file_path = ''
        self.excel_row = -1


def is_region_overlap(x0, y0, x1, y1, x2, y2, x3, y3):
    """If two regions are overlap.

    Arguments:
        x0: Float, region1.x0
        y0: Float, region1.y0
        x1: Float, region1.x1
        y1: Float, region1.y1
        x2: Float, region2.x0
        y2: Float, region2.y0
        x3: Float, region2.x1
        y3: Float, region2.y1

    Returns:
        Boolen, True-overlapped, False-not overlapped.
    """


    overlap_flag = True

    minx = max(x0 , x2)
    miny = max(y0 , y2)
    maxx = min(x1 , x3)
    maxy = min(y1 , y3)

    if minx > maxx or miny > maxy:
        overlap_flag = False

    return overlap_flag


def is_long_coordinate(text):
    digit_count = 0
    for t in text:
        if t.isdigit():
            digit_count += 1

    if float(digit_count / len(text)) > 0.5:
        return True
    else:
        return False


def get_region_type(text):
    region_type = -1
    if text.find('部') > 0 and (not text.startswith('（')) and len(text) > 1 and (text[text.find('部')+1] == '\n' or text[text.find('部')+1] == ' ' or text[text.find('部')+1] == '.') and text[text.find('部')-1].isdigit():
        region_type = 1
    elif text.find('章') > 0 and (not text.startswith('（')) and len(text) > 1 and (text[text.find('章')+1] == '\n' or text[text.find('章')+1] == ' ' or text[text.find('章')+1] == '.') and text[text.find('章')-1].isdigit():
        region_type = 2
    elif text.find('節') > 0 and (not text.startswith('（')) and len(text) > 1 and (text[text.find('節')+1] == '\n' or text[text.find('節')+1] == ' ' or text[text.find('節')+1] == '.') and text[text.find('節')-1].isdigit():
        region_type = 3
    elif text[0].isdigit() and len(text) > 1 and text[1] == '．':
        region_type = 4
    elif text.find('第') >-1 and text.startswith('【') and len(text) > 1 and text[text.find('第')+1].isdigit():
        region_type = 6
    elif len(text) < 4:
        text = text.replace('\n', '')
        if text.isdigit():
            region_type = 7
    else:
        region_type = 5

    return region_type

def get_region_num(region_type, text):
    if region_type in [1, 2, 3]:
        for i in range(text.find('第')+1, len(text)):
            if not text[i].isdigit():
                if text[text.find('第')+1:i].isdigit():
                    try:
                        return int(text[text.find('第')+1:i])
                    except:
                        return -1
    elif region_type == 4:
        for i in range(len(text)):
            if not text[i].isdigit():
                return int(text[:i])

    return -1

def get_title_num(text):
    for i in range(text.find('第')+1, len(text)):
        if (not text[i].isdigit()) and (text[i] != '-'):
            try:
                return text[text.find('第')+1:i]
            except:
                return ''                

def is_related_paragraph(title_num, text):
    """If the paragraph is related to the title number.

    """


    if title_num in text:
        return True
    else:
        return False


def has_numbers(input_string):
    return bool(re.search(r'\d', input_string))


def has_special_digits(input_string):
    detected = False
    special_digits = ['①', '②', '③', '④', '⑤', '⑥', '⑦', '⑧', '⑨', '⑩']
    for special_digit in special_digits:
        if input_string.startswith(special_digit):
            detected = True
            break

    special_digits = ['ア', 'イ', 'ウ', 'エ', 'オ', 'カ', 'キ', 'ク', 'ケ', 'コ']
    if input_string.startswith('（'):
        for special_digit in special_digits:
            if input_string.find(special_digit) >-1:
                detected = True
                break

    if input_string.startswith('（') and input_string[input_string.find('（')+1].isdigit():
         detected = True

    return detected


def text_output1(pdf_path, output_path):
    # 新建文件
    #wb = openpyxl.Workbook() 
    wb = openpyxl.load_workbook(output_path)

    fp = open(pdf_path, 'rb')  # 以二进制读模式打开
    # 用文件对象来创建一个pdf文档分析器
    parser = PDFParser(fp)
    # 创建一个PDF文档
    doc = PDFDocument()
    # 连接分析器 与文档对象
    parser.set_document(doc)
    doc.set_parser(parser)

    # 提供初始化密码
    # 如果没有密码 就创建一个空的字符串
    doc.initialize()

    # 检测文档是否提供txt转换，不提供就忽略
    if not doc.is_extractable:
        raise PDFTextExtractionNotAllowed
    else:
        # 创建PDf 资源管理器 来管理共享资源
        rsrcmgr = PDFResourceManager()
        # 创建一个PDF设备对象
        laparams = LAParams()
        device = PDFPageAggregator(rsrcmgr, laparams=laparams)
        # 创建一个PDF解释器对象
        interpreter = PDFPageInterpreter(rsrcmgr, device)

        # 用来计数页面，图片，曲线，figure，水平文本框等对象的数量
        num_page, num_image, num_curve, num_figure, num_TextBoxHorizontal = 0, 0, 0, 0, 0

        # 循环遍历列表，每次处理一个page的内容
        for page in doc.get_pages(): # doc.get_pages() 获取page列表
            num_page += 1  # 页面增一
            text_line = ''
            interpreter.process_page(page)
            # 接受该页面的LTPage对象
            layout = device.get_result()
            row = 1
            for x in layout:
                if isinstance(x,LTImage):  # 图片对象
                    num_image += 1
                if isinstance(x,LTCurve):  # 曲线对象
                    num_curve += 1
                if isinstance(x,LTFigure):  # figure对象
                    num_figure += 1
                if isinstance(x, LTTextBoxHorizontal):  # 获取文本内容
                    num_TextBoxHorizontal += 1  # 水平文本框对象增一
                    # 保存文本内容
                    text_line += x.get_text()
                    if text_line.endswith('。'):
                        text_line += ' '

            text_line.replace('\n', '')
            ws = wb.create_sheet()
            ws.title = str(num_page)
            ws.column_dimensions['A'].width = 50
            cell_name = 'A1'
            row += 1
            try:
                ws[cell_name] = text_line
                alignment = Alignment(wrap_text=True)
                ws[cell_name].alignment = alignment
            except openpyxl.utils.exceptions.IllegalCharacterError:
                continue

    # 写入文件
    #sheet = workbook.active
    #sheet['A1']='data'
    # 保存文件 
    wb.save(output_path)

def img_output2(pdf_path, output_path):
    
    if os.path.exists(output_path):
        wb = openpyxl.load_workbook(output_path)
    else:
        wb = openpyxl.Workbook() 

    doc = fitz.open(pdf_path)
    for i in range(len(doc)):
        img_index = 1
        for img in doc.getPageImageList(i):
            xref = img[0]
            pix = fitz.Pixmap(doc, xref)
            if pix.width < 100 or pix.height < 100:
                pix = None
                continue
            print('Page: %d, Xref: %d' % (i, xref))
            if pix.n < 5:       # this is GRAY or RGB
                try:
                    img_name = IMG_PATH+"p%s-%s.png" % (i, xref)
                    pix.writePNG(img_name)
                    img = Image(img_name)
                    cell_name = 'B' + str(img_index)
                    ws = wb[str(i)]
                    ws.add_image(img, cell_name) 
                    img_index += 1
                    print('Add image: %s to Page: %s in cell: %s' % (img_name, str(i), cell_name))
                except RuntimeError:
                    pix = None
                    continue
            else:               # CMYK: convert to RGB first
                img_name = IMG_PATH+"p%s-%s.png" % (i, xref)
                pix1 = fitz.Pixmap(fitz.csRGB, pix)
                pix1.writePNG(img_name)
                img = Image(img_name)
                cell_name = 'B' + str(img_index)
                ws = wb[str(i)]
                ws.add_image(img, cell_name) 
                img_index += 1
                print('Add image: %s to Page: %s in cell: %s' % (img_name, str(i), cell_name))
                pix1 = None
            pix = None

    wb.save('test.xlsx')


def translate_points(x0, y0, x1, y1, jpg_img, pdf_height):
    jpg_height = jpg_img.shape[0]
    scale = float(jpg_height / pdf_height)

    return x0 * scale, jpg_height - y1 * scale, x1 * scale, jpg_height - y0 * scale

def get_eucldist(x0, y0, x1, y1):
    a = np.array((x0 ,y0))
    b = np.array((x1, y1))

    return np.linalg.norm(a-b)


def get_rect_dist(x1, y1, x1b, y1b, x2, y2, x2b, y2b):
    left = x2b < x1
    right = x1b < x2
    bottom = y2b < y1
    top = y1b < y2

    if top and left:
        return get_eucldist(x1, y1b, x2b, y2)
    elif left and bottom:
        return get_eucldist(x1, y1, x2b, y2b)
    elif bottom and right:
        return get_eucldist(x1b, y1, x2, y2b)
    elif right and top:
        return get_eucldist(x1b, y1b, x2, y2)
    elif left:
        return x1 - x2b
    elif right:
        return x2 - x1b
    elif bottom:
        return y1 - y2b
    elif top:
        return y2 - y1b
    else: # rectangles intersect
        return 0.

def text_output(pdf_path, output_path):
    # 新建文件
    if os.path.exists(output_path): 
        wb = openpyxl.load_workbook(output_path)
    else:
        wb = openpyxl.Workbook()

    with open(pdf_path, 'rb') as temp_fp:
        pdf_height = PdfFileReader(temp_fp).getPage(0).mediaBox[3]
        pdf_width = PdfFileReader(temp_fp).getPage(0).mediaBox[2]
        
    print(pdf_height)
    print(pdf_width)

    fp = open(pdf_path, 'rb')  # 以二进制读模式打开
    # 用文件对象来创建一个pdf文档分析器
    parser = PDFParser(fp)
    # 创建一个PDF文档
    doc = PDFDocument()
    # 连接分析器 与文档对象
    parser.set_document(doc)
    doc.set_parser(parser)

    # 提供初始化密码
    # 如果没有密码 就创建一个空的字符串
    doc.initialize()

    # 检测文档是否提供txt转换，不提供就忽略
    if not doc.is_extractable:
        raise PDFTextExtractionNotAllowed
    else:
        # 创建PDf 资源管理器 来管理共享资源
        rsrcmgr = PDFResourceManager()
        # 创建一个PDF设备对象
        laparams = LAParams()
        device = PDFPageAggregator(rsrcmgr, laparams=laparams)
        # 创建一个PDF解释器对象
        interpreter = PDFPageInterpreter(rsrcmgr, device)

        # 用来计数页面，图片，曲线，figure，水平文本框等对象的数量
        num_page, num_image, num_curve, num_figure, num_TextBoxHorizontal = 0, 0, 0, 0, 0

        ws = wb.active
        ws['A1'] = 'No.'
        ws['B1'] = 'タイプ'
        ws['C1'] = '頁番号'
        ws['D1'] = '部'
        ws['E1'] = '章'
        ws['F1'] = '節'
        ws['G1'] = '小節'
        ws['H1'] = '段落'
        ws['I1'] = '本文保存ファイル'
        ws['J1'] = '図保存ファイル'
        ws['K1'] = '関連文ファイル'

        # 循环遍历列表，每次处理一个page的内容
        row = 2
        part_num = -1
        chapter_num = -1
        section_num = -1
        sub_section_num = -1
        paragraph_text = ''
        for page in doc.get_pages(): # doc.get_pages() 获取page列表
            num_page += 1  # 页面增一
            print(num_page)
            if num_page < 11:
                continue

            if num_page == 187:
                ccc = 0

            interpreter.process_page(page)
            # 接受该页面的LTPage对象
            layout = device.get_result()
            region_list = []
            page_num = -1
            paragraph_num = 1
            jpg_file_path = '%sjpg/pg_%04d.jpg' % (TEMP_PATH, num_page)
            jpg_img = cv2.imread(jpg_file_path)
            title_region_list = []
            # Get title region
            for x in layout:
                conv_x0, conv_y0, conv_x1, conv_y1 = translate_points(x.x0, x.y0, x.x1, x.y1, jpg_img, pdf_height)

                if isinstance(x, LTTextBoxHorizontal):
                    cur_text = x.get_text()
                    #print(cur_text)
                    region_type = get_region_type(cur_text)
                    if region_type != 6:
                        continue

                    # Title
                    title_num = get_title_num(cur_text)
                    text_region = TextRegion(region_type, conv_x0, conv_y0, conv_x1, conv_y1, x._avg_lineheight, cur_text)
                    text_region.part_number = part_num
                    text_region.chapter_number = chapter_num
                    text_region.section_number = section_num
                    text_region.sub_section_number = sub_section_num
                    text_region.title_number = title_num
                    title_region_list.append(text_region)
            
            img_corners = []
            # Image Prosessing
            img_file_path = '%spng/pg_%04d.png' % (TEMP_PATH, num_page)
            corner_pixels = eifp.find_contours(img_file_path, False, False)
            page_image = cv2.imread(img_file_path)
            img_num_list = []
            for corner_pixel in corner_pixels:
                min_dist = 999999
                img_num = ''
                for region in title_region_list:
                    dist = np.sqrt(np.square(corner_pixel[1]-region.x0) + np.square(corner_pixel[0]-region.y0)) 
                    #dist = abs(corner_pixel[1]-region.x0)
                    if dist < min_dist and (region.title_number not in img_num_list):
                        min_dist = dist
                        img_num = region.title_number

                if len(img_num) > 0:
                    img_corners.append([corner_pixel[1], corner_pixel[0], corner_pixel[3], corner_pixel[2]])
                    img_num_list.append(img_num)
                    img_file_dir = DATA_PATH + 'IMG/'
                    if not os.path.exists(img_file_dir):
                        os.makedirs(img_file_dir)
                    img_output_name = img_file_dir + img_num + '.png'
                    img_region = page_image[corner_pixel[0]:corner_pixel[2], corner_pixel[1]:corner_pixel[3], :]
                    cv2.rectangle(jpg_img, (int(corner_pixel[1]),int(corner_pixel[0])), (int(corner_pixel[3]),int(corner_pixel[2])), (255,0,0), 1)
                    cv2.imwrite(img_output_name, img_region)
                    print('%s saved' % img_output_name)
                    for region in title_region_list:
                        if region.title_number == img_num:
                            region.img_path = img_output_name    

            region_count = 0
            for x in layout:
                conv_x0, conv_y0, conv_x1, conv_y1 = translate_points(x.x0, x.y0, x.x1, x.y1, jpg_img, pdf_height)

                if isinstance(x, LTTextBoxHorizontal):  # 获取文本内容                
                    # 保存文本内容
                    cur_text = x.get_text().lstrip()
                    #print(cur_text)
                    region_type = get_region_type(cur_text)
                    is_img_text = False
                    for img_corner in img_corners:
                        if is_region_overlap(conv_x0, conv_y0, conv_x1, conv_y1, img_corner[0], img_corner[1], img_corner[2], img_corner[3]) or get_rect_dist(conv_x0, conv_y0, conv_x1, conv_y1, img_corner[0], img_corner[1], img_corner[2], img_corner[3]) < 15:
                            is_img_text = True
                            break

                    if is_img_text:
                        continue

                    """
                    if cur_text.startswith('出典'):
                        continue
                    """

                    if cur_text.find('出典') >-1:
                        continue

                    if cur_text.find('注') > -1:
                        if cur_text[cur_text.find('注')+1].isdigit():
                            continue

                    

                    # 1: Part, 2: Chapter, 3: Section, 4: Sub-section, 5: Paragraph, 6: Title, 7: Page number
                    if region_type == 1:
                        # Part
                        part_num = get_region_num(region_type, cur_text)
                        chapter_num = -1
                        section_num = -1
                        sub_section_num = -1
                        text_region = TextRegion(region_type, conv_x0, conv_y0, conv_x1, conv_y1, x._avg_lineheight, cur_text)                        
                        text_region.part_number = part_num
                        region_list.append(text_region)
                        cv2.rectangle(jpg_img, (int(conv_x0),int(conv_y0)), (int(conv_x1),int(conv_y1)), (0,255,0), 1)
                        font = cv2.FONT_HERSHEY_SIMPLEX
                        cv2.putText(jpg_img, str(region_count), (int(conv_x0), int(conv_y0)), font, 1, (204, 0, 0), 2)
                        region_count += 1
                    elif region_type == 2:
                        # Chapter
                        chapter_num = get_region_num(region_type, cur_text)
                        section_num = -1
                        sub_section_num = -1
                        text_region = TextRegion(region_type, conv_x0, conv_y0, conv_x1, conv_y1, x._avg_lineheight, cur_text)
                        text_region.part_number = part_num
                        text_region.chapter_number = chapter_num
                        region_list.append(text_region)
                        cv2.rectangle(jpg_img, (int(conv_x0),int(conv_y0)), (int(conv_x1),int(conv_y1)), (0,255,0), 1)
                        font = cv2.FONT_HERSHEY_SIMPLEX
                        cv2.putText(jpg_img, str(region_count), (int(conv_x0), int(conv_y0)), font, 1, (204, 0, 0), 2)
                        region_count += 1
                    elif region_type == 3:
                        # Section
                        section_num = get_region_num(region_type, cur_text)
                        sub_section_num = -1
                        text_region = TextRegion(region_type, conv_x0, conv_y0, conv_x1, conv_y1, x._avg_lineheight, cur_text)
                        text_region.part_number = part_num
                        text_region.chapter_number = chapter_num
                        text_region.section_number = section_num
                        region_list.append(text_region)
                        cv2.rectangle(jpg_img, (int(conv_x0),int(conv_y0)), (int(conv_x1),int(conv_y1)), (0,255,0), 1)
                        font = cv2.FONT_HERSHEY_SIMPLEX
                        cv2.putText(jpg_img, str(region_count), (int(conv_x0), int(conv_y0)), font, 1, (204, 0, 0), 2)
                        region_count += 1
                    elif region_type == 4:
                        # Sub-section
                        sub_section_num = get_region_num(region_type, cur_text)
                        text_region = TextRegion(region_type, conv_x0, conv_y0, conv_x1, conv_y1, x._avg_lineheight, cur_text)
                        text_region.part_number = part_num
                        text_region.chapter_number = chapter_num
                        text_region.section_number = section_num
                        text_region.sub_section_number = sub_section_num
                        region_list.append(text_region)
                        cv2.rectangle(jpg_img, (int(conv_x0),int(conv_y0)), (int(conv_x1),int(conv_y1)), (0,255,0), 1)
                        font = cv2.FONT_HERSHEY_SIMPLEX
                        cv2.putText(jpg_img, str(region_count), (int(conv_x0), int(conv_y0)), font, 1, (204, 0, 0), 2)
                        region_count += 1
                    elif region_type == 5:
                        if len(cur_text) < 5 or abs(x.x0 - x.x1) < pdf_width / 10:
                            continue

                        if x.y0 < pdf_height / 6 and x.y1 < pdf_height / 6 and cur_text[0].isdigit():
                            continue

                        if is_long_coordinate(cur_text):
                            continue

                        if abs(x.y0 - x.y1) < 10:
                            continue

                        if x.x0 > pdf_width or x.x1 > pdf_width or (x.y0 > 0.9 * float(pdf_height) and x.y1 > 0.9 * float(pdf_height)):
                            continue

                        if has_special_digits(cur_text):
                            continue
                            
                        if x.width > pdf_width / 2:
                            continue

                        cv2.rectangle(jpg_img, (int(conv_x0),int(conv_y0)), (int(conv_x1),int(conv_y1)), (0,255,0), 1)
                        font = cv2.FONT_HERSHEY_SIMPLEX
                        cv2.putText(jpg_img, str(region_count), (int(conv_x0), int(conv_y0)), font, 1, (204, 0, 0), 2)
                        region_count += 1
                        paragraph_text += cur_text
                        texts = paragraph_text.split('。\n')
                        """
                        texts = []
                        for temp_text in temp_texts:
                            texts.append(temp_text.split('）\n'))
                        """
                        if paragraph_text[-2] == '。' or paragraph_text[-2] == '）':
                            for text_index, sp_text in enumerate(texts):
                                if len(sp_text) < 1:
                                    continue
                                sp_text = sp_text.replace('\n', '')
                                text_region = TextRegion(region_type, conv_x0, conv_y0, conv_x1, conv_y1, x._avg_lineheight, sp_text)
                                text_region.part_number = part_num
                                text_region.chapter_number = chapter_num
                                text_region.section_number = section_num
                                text_region.sub_section_number = sub_section_num                                
                                text_region.paragraph_number = paragraph_num
                                title_num = get_title_num(cur_text)
                                if len(title_num) > 0:
                                    text_region.title_number = title_num
                                region_list.append(text_region)
                                paragraph_num += 1

                            paragraph_text = ''
                        else:
                            for text_index, sp_text in enumerate(texts):
                                if text_index != len(texts) - 1:
                                    if len(sp_text) < 1:
                                        continue
                                    sp_text = sp_text.replace('\n', '')
                                    text_region = TextRegion(region_type, conv_x0, conv_y0, conv_x1, conv_y1, x._avg_lineheight, sp_text)
                                    text_region.part_number = part_num
                                    text_region.chapter_number = chapter_num
                                    text_region.section_number = section_num
                                    text_region.sub_section_number = sub_section_num
                                    text_region.paragraph_number = paragraph_num
                                    title_num = get_title_num(cur_text)
                                    if len(title_num) > 0:
                                        text_region.title_number = title_num
                                    region_list.append(text_region)
                                    paragraph_num += 1
                                else:
                                    paragraph_text = sp_text
                    
                    elif region_type == 6:
                        # Title
                        """
                        title_num = get_title_num(cur_text)
                        text_region = TextRegion(region_type, conv_x0, conv_y0, conv_x1, conv_y1, x._avg_lineheight, cur_text)
                        text_region.part_number = part_num
                        text_region.chapter_number = chapter_num
                        text_region.section_number = section_num
                        text_region.sub_section_number = sub_section_num
                        text_region.title_number = title_num
                        region_list.append(text_region)
                        """
                        cv2.rectangle(jpg_img, (int(conv_x0),int(conv_y0)), (int(conv_x1),int(conv_y1)), (0,255,0), 1)
                        font = cv2.FONT_HERSHEY_SIMPLEX
                        cv2.putText(jpg_img, str(region_count), (int(conv_x0), int(conv_y0)), font, 1, (204, 0, 0), 2)
                        region_count += 1
                    elif region_type == 7:
                        page_num = int(cur_text)
                        cv2.rectangle(jpg_img, (int(conv_x0),int(conv_y0)), (int(conv_x1),int(conv_y1)), (0,255,0), 1)
                        font = cv2.FONT_HERSHEY_SIMPLEX
                        cv2.putText(jpg_img, str(region_count), (int(conv_x0), int(conv_y0)), font, 1, (204, 0, 0), 2)
                        region_count += 1                        

            for title_region in title_region_list:
                region_list.append(title_region)

            cv2.imwrite('%sseg_jpg/pg_%04d.jpg' % (TEMP_PATH, num_page), jpg_img)

            #ws = wb.create_sheet()
            #ws.title = str(num_page)
            img_title_list = {}
            for region in region_list:            
                if region.region_type == 6:
                    if region.title_number in img_title_list.values():
                        continue
                    
                # Region No.
                region.excel_row = row
                cell_name = 'A' + str(row)
                try:
                    ws[cell_name] = str(row - 1)
                    alignment = Alignment(wrap_text=True)
                    ws[cell_name].alignment = alignment
                except openpyxl.utils.exceptions.IllegalCharacterError:
                    pass

                # Region type
                cell_name = 'B' + str(row)
                try:
                    ws[cell_name] = dic_region_type[region.region_type]
                    alignment = Alignment(wrap_text=True)
                    ws[cell_name].alignment = alignment
                except openpyxl.utils.exceptions.IllegalCharacterError:
                    pass

                # Page Number
                cell_name = 'C' + str(row)
                try:
                    ws[cell_name] = str(num_page-4)
                    alignment = Alignment(wrap_text=True)
                    ws[cell_name].alignment = alignment
                except openpyxl.utils.exceptions.IllegalCharacterError:
                    pass
                
                # Part Number
                cell_name = 'D' + str(row)
                try:
                    ws[cell_name] = str(region.part_number)
                    alignment = Alignment(wrap_text=True)
                    ws[cell_name].alignment = alignment
                except openpyxl.utils.exceptions.IllegalCharacterError:
                    pass

                if region.region_type in [2, 3, 4, 5, 6]:
                    # Chapter Number
                    cell_name = 'E' + str(row)
                    try:
                        ws[cell_name] = str(region.chapter_number)
                        alignment = Alignment(wrap_text=True)
                        ws[cell_name].alignment = alignment
                    except openpyxl.utils.exceptions.IllegalCharacterError:
                        pass

                if region.region_type in [3, 4, 5, 6]:
                    # Section Number
                    cell_name = 'F' + str(row)
                    try:
                        ws[cell_name] = str(region.section_number)
                        alignment = Alignment(wrap_text=True)
                        ws[cell_name].alignment = alignment
                    except openpyxl.utils.exceptions.IllegalCharacterError:
                        pass

                if region.region_type in [4, 5, 6]:
                    # Sub_section Number
                    cell_name = 'G' + str(row)
                    try:
                        ws[cell_name] = str(region.sub_section_number)
                        alignment = Alignment(wrap_text=True)
                        ws[cell_name].alignment = alignment
                    except openpyxl.utils.exceptions.IllegalCharacterError:
                        pass

                if region.region_type == 5:
                    # Paragraph Number
                    cell_name = 'H' + str(row)
                    try:
                        ws[cell_name] = str(region.paragraph_number)
                        alignment = Alignment(wrap_text=True)
                        ws[cell_name].alignment = alignment
                    except openpyxl.utils.exceptions.IllegalCharacterError:
                        pass
                # Text
                cell_name = 'I' + str(row)

                try:
                    file_dir = DATA_PATH + 'TXT/'
                    if not os.path.exists(file_dir):
                        os.makedirs(file_dir)
                    file_path = file_dir + str(row - 1) + '.txt'
                    with open(file_path, 'w') as f:
                        f.write(region.text)

                    region.text_file_path = file_path 

                    ws[cell_name] = str(file_path)
                    alignment = Alignment(wrap_text=True)
                    ws[cell_name].alignment = alignment
                except openpyxl.utils.exceptions.IllegalCharacterError:
                    pass

                if region.region_type == 6:
                    # Title Number
                    '''
                    cell_name = 'J' + str(row)
                    try:
                        ws[cell_name] = region.title_number
                        alignment = Alignment(wrap_text=True)
                        ws[cell_name].alignment = alignment
                    except openpyxl.utils.exceptions.IllegalCharacterError:
                        pass
                    '''

                    img_title_list[row] = region.title_number

                    # Img Path
                    cell_name = 'J' + str(row)
                    try:
                        ws[cell_name] = region.img_path
                        alignment = Alignment(wrap_text=True)
                        ws[cell_name].alignment = alignment
                    except openpyxl.utils.exceptions.IllegalCharacterError:
                        pass
                
                row += 1

                if region.region_type == 1:
                    with open('part.txt', 'a') as f:
                        f.write(region.text)
                        f.write('\n')
                elif region.region_type == 2:
                    with open('chapter.txt', 'a') as f:
                        f.write(region.text)
                        f.write('\n')
                elif region.region_type == 3:
                    with open('section.txt', 'a') as f:
                        f.write(region.text)
                        f.write('\n')
                elif region.region_type == 4:
                    with open('sub_section.txt', 'a') as f:
                        f.write(region.text)
                        f.write('\n')
                elif region.region_type == 6:
                    with open('title.txt', 'a') as f:
                        f.write(region.text)
                        f.write('\n')

            for key, value in img_title_list.items():
                for region in region_list:
                    if region.region_type != 5:
                        continue
                        
                    if is_related_paragraph(value, region.text):
                        cell_name = 'K' + str(key)
                        try:
                            ws[cell_name] = region.text_file_path
                            alignment = Alignment(wrap_text=True)
                            ws[cell_name].alignment = alignment
                        except openpyxl.utils.exceptions.IllegalCharacterError:
                            pass
            '''
            try:
                ws = wb[str(num_page)]
            except:
                ws = wb.create_sheet()
                ws.title = str(num_page)
            ws.column_dimensions['A'].width = 50

            for i in range(len(region_list)):
                region = region_list[i]
                cell_name = 'A' + str(i+1)
                try:
                    ws[cell_name] = region.text
                    alignment = Alignment(wrap_text=True)
                    ws[cell_name].alignment = alignment
                except openpyxl.utils.exceptions.IllegalCharacterError:
                    continue

                cell_name = 'B' + str(i+1)
                try:
                    ws[cell_name] = region.region_type
                    alignment = Alignment(wrap_text=True)
                    ws[cell_name].alignment = alignment
                except openpyxl.utils.exceptions.IllegalCharacterError:
                    continue
            '''

    wb.save(output_path)

def img_output(pdf_path, output_path):
    
    wb  =  xlsxwriter.Workbook(output_path)

    doc = fitz.open(pdf_path)
    for i in range(len(doc)):
        img_index = 1
        ws = wb.add_worksheet(str(i+1))
        im_list = doc.getPageImageList(i)
        for img in doc.getPageImageList(i):
            xref = img[0]
            pix = fitz.Pixmap(doc, xref)
            if pix.width < 100 or pix.height < 100:
                pix = None
                continue
            print('Page: %d, Xref: %d' % (i, xref))
            if pix.n < 5:       # this is GRAY or RGB
                try:
                    img_name = IMG_PATH+"p%s-%s.png" % (i, xref)
                    pix.writePNG(img_name)
                    img = Image(img_name)
                    cell_name = 'B' + str(img_index*10)
                    ws.insert_image(cell_name, img_name, {'positioning': 1})
                    img_index += 1
                    print('Add image: %s to Page: %s in cell: %s' % (img_name, str(i), cell_name))
                except RuntimeError:
                    pix = None
                    continue
            else:               # CMYK: convert to RGB first
                img_name = IMG_PATH+"p%s-%s.png" % (i, xref)
                pix1 = fitz.Pixmap(fitz.csRGB, pix)
                pix1.writePNG(img_name)
                img = Image(img_name)
                cell_name = 'B' + str(img_index*10)
                ws.insert_image(cell_name, img_name, {'positioning': 1})
                img_index += 1
                print('Add image: %s to Page: %s in cell: %s' % (img_name, str(i), cell_name))
                pix1 = None
            pix = None

    wb.close()


def slide_output(ledger_path, output_path):
    """Output images and related text to slide
    
    """


    # Read ledge
    wb = openpyxl.load_workbook(ledger_path)
    ws = wb.active

    # Slide init
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    #title_num_list = []
    for row in range(2, ws.max_row+1):
        page_num = int(ws['C'+str(row)].value)
        region_type = ws['B'+str(row)].value
        if len(region_type) < 1:
            break
        if region_type == '図タイトル':
            title_path = ws['I'+str(row)].value
            img_path = ws['J'+str(row)].value
            relate_text_path = ws['K'+str(row)].value
            title_num = get_title_num(ta.read_text(title_path))
            #if title_num in title_num_list:
            #    continue
            #else:
            #    title_num_list.append(title_num)

            if img_path is None or relate_text_path is None:
                continue

            # Create slide   
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Page number
            txBox = slide.shapes.add_textbox(1, 1, 1, 1)
            tf = txBox.text_frame
            tf.text = 'Page: %d' % page_num

            # Image
            pic = slide.shapes.add_picture(img_path, Inches(0.5), Inches(1), width=Inches(4))

            # Image title
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(4), Inches(1))
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = '図タイトル:'
            p.font.bold = True

            text = ta.read_text(title_path)
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(4), Inches(1))
            tf = txBox.text_frame
            tf.text = text
            tf.margin_bottom = Inches(0.08)
            tf.margin_left = 0
            tf.vertical_anchor = MSO_ANCHOR.TOP
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

            # Related text
            txBox = slide.shapes.add_textbox(Inches(5), Inches(0.5), Inches(5), Inches(1))
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = '関連文:'
            p.font.bold = True

            text_paths = relate_text_path.split('\n')
            left = Inches(5)
            top = Inches(1)
            width = Inches(5)
            height = Inches(1)
            for text_path in text_paths:
                if len(text_path) == 0:
                    continue

                text = ta.read_text(text_path)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                top += height
                tf = txBox.text_frame
                #tf.text = text
                tf.margin_bottom = Inches(0.08)
                tf.margin_left = 0
                tf.vertical_anchor = MSO_ANCHOR.TOP
                tf.word_wrap = True
                tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                p = tf.add_paragraph()
                p.text = text
                p.font.size = Pt(12)

            """
            text = ta.read_text(relate_text_path)
            txBox = slide.shapes.add_textbox(Inches(5), Inches(1), Inches(5), Inches(1))
            tf = txBox.text_frame
            #tf.text = text
            tf.margin_bottom = Inches(0.08)
            tf.margin_left = 0
            tf.vertical_anchor = MSO_ANCHOR.TOP
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            p = tf.add_paragraph()
            p.text = text
            p.font.size = Pt(12)
            #p = tf.add_paragraph()
            #.text = text
            """

    prs.save(output_path)


def check_related_text(ledger_path, output_path):

    # Read ledge
    wb = openpyxl.load_workbook(ledger_path)
    ws = wb.active


    for row in range(2, ws.max_row+1):
        print(row)
        region_type = ws['B'+str(row)].value
        if len(region_type) < 1:
            break
        if region_type == '図タイトル':
            title_path = ws['I'+str(row)].value

            if title_path is None:
                continue

            text = ta.read_text(title_path)
            title_num = get_title_num(text)
            related_texts = ''
            for para_row in range(row-15, row+15):
                if para_row < 2:
                    para_row = 2

                if ws['B'+str(para_row)].value != '段落':
                    continue

                para_text = ta.read_text(ws['I'+str(para_row)].value)
                if is_related_paragraph(title_num, para_text):
                    related_texts += str(ws['I'+str(para_row)].value)
                    related_texts += '\n'
                    """
                    cell_name = 'K' + str(row)
                    try:
                        ws[cell_name] = ws['I'+str(para_row)].value
                        alignment = Alignment(wrap_text=True)
                        ws[cell_name].alignment = alignment
                    except openpyxl.utils.exceptions.IllegalCharacterError:
                        pass
                    """
            
            cell_name = 'K' + str(row)
            try:
                ws[cell_name] = related_texts
                alignment = Alignment(wrap_text=True)
                ws[cell_name].alignment = alignment
            except openpyxl.utils.exceptions.IllegalCharacterError:
                pass

    wb.save(output_path)


def main():
    pdf_path = DATA_PATH + 'whitepaper2018pdf_all.pdf'
    ledger_path = 'energy.xlsx'
    slide_path = 'energy_slide.pptx'
    text_output(pdf_path, ledger_path)
    check_related_text(ledger_path, 'energy_c.xlsx')
    slide_output('energy_c.xlsx', slide_path)


if __name__ == '__main__':
    main()