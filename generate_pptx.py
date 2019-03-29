# -*- coding: utf-8 -*-

import glob
import pickle
from datetime import datetime
from pptx import Presentation
from pptx.util import Pt
import MeCab

m_ocha = MeCab.Tagger(r'-Ochasen -d C:\Users\hori\workspace\encoder-decoder-sentence-chainer-master\mecab-ipadic-neologd')

path = "dictionary/wakati_filtered"
dict_path_list = glob.glob(path + "/*")

today = datetime.today().strftime("%Y%m%d")

prs = Presentation()

title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "アーカイブ頻出単語"
subtitle.text = str(today) + "堀智之"

for dict_path in dict_path_list:
    fname = dict_path.replace(path, "").replace("\\", "")[:-7]
    dictionary = pickle.load(open(dict_path, "rb"))
    sorted_word_list = []
    sorted_num_list = []
    sorted_category_list = []
    for k, v in sorted(dictionary.items(), key=lambda x: -x[1]):
        sorted_word_list.append(k)
        sorted_num_list.append(v)
        node = m_ocha.parseToNode(k)
        node = node.next
        fields = node.feature.split(",")
        sorted_category_list.append(fields[0])
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "ファイル名" + fname
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)  # font size
    title_shape.text_frame.paragraphs[0].font.bold = True  # font bold
    
    tf = body_shape.text_frame
    tf.text = "頻出単語、品詞、出現回数"
    tf.paragraphs[0].font.size = Pt(18)  # font size
    tf.paragraphs[0].font.bold = True  # font bold
    for i in range(15):
        p = tf.add_paragraph()
        p.text = "「" + str(sorted_word_list[i]) + "(" + sorted_category_list[i] + ")」・・・" + str(sorted_num_list[i])
        p.level = 1
        p.font.size = Pt(14)  # font size
    
prs.save("result_meishi%s.pptx" % today)