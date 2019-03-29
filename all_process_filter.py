# -*- coding: utf-8 -*-

from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from pdfminer.converter import TextConverter
from pdfminer.layout import LAParams
from pdfminer.pdfpage import PDFPage
import glob
import MeCab
import codecs
from datetime import datetime
from pptx import Presentation
from pptx.util import Pt
from title_master import title_master
from eliminate_word_master import eliminate_word_master

title_master_class = title_master()
title_dict = title_master_class.get_dict()
eliminate_word_list = eliminate_word_master()
eliminate_word_list.load_list("eliminate_word_list.pickle")

m = MeCab.Tagger(r'-Owakati -d C:\Users\hori\workspace\encoder-decoder-sentence-chainer-master\mecab-ipadic-neologd')
m_ocha = MeCab.Tagger(r'-Ochasen -d C:\Users\hori\workspace\encoder-decoder-sentence-chainer-master\mecab-ipadic-neologd')
"""
file_list = glob.glob('pdf/*')
for input_path in file_list:
    output_path = 'text/' + input_path[4:-4] + '.txt'
    rsrcmgr = PDFResourceManager()
    codec = 'utf-8'
    params = LAParams()
    with open(output_path, "wb") as output:
        device = TextConverter(rsrcmgr, output, codec=codec, laparams=params)
        with open(input_path, 'rb') as input:
            interpreter = PDFPageInterpreter(rsrcmgr, device)
            for page in PDFPage.get_pages(input):
                interpreter.process_page(page)
        device.close()
"""
text_path_list = glob.glob('text/*.txt')

today = datetime.today().strftime("%Y%m%d")

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "アーカイブ頻出単語"
subtitle.text = str(today) + "堀智之"

for text_path in text_path_list:
    print(text_path)
    f = codecs.open(text_path, 'r', 'utf-8')
    text = f.read()
    text = text.replace(' ', '')
    text = text.replace('\n', '')
    
    node = m_ocha.parseToNode(text)
    
    word_count_dict = {}
    word_category_dict = {}
    while node:
        fields = node.feature.split(",")
        word = node.surface
        
        
        if fields[0] == '名詞' or fields[0] == '動詞' or fields[0] == '形容詞' or fields[0] == '形容動詞':
            if not eliminate_word_list.is_include(word) and fields[1] != '数':
                if word in word_count_dict:
                    count = word_count_dict[word]
                    word_count_dict[word] = count + 1
                else:
                    word_count_dict[word] = 1
                if word not in word_category_dict:
                    word_category_dict[word] = fields[0]
        
        """
        if fields[0] == '名詞':
            if not eliminate_word_list.is_include(word) and fields[1] != '数':
                if word in word_count_dict:
                    count = word_count_dict[word]
                    word_count_dict[word] = count + 1
                else:
                    word_count_dict[word] = 1
                if word not in word_category_dict:
                    word_category_dict[word] = fields[0]
        """
        
        node = node.next
        
    fname = text_path[5:-4]
    sorted_word_list = []
    sorted_num_list = []
    sorted_category_list = []
    for k, v in sorted(word_count_dict.items(), key=lambda x: -x[1]):
        sorted_word_list.append(k)
        sorted_num_list.append(v)
        sorted_category_list.append(word_category_dict[k])
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    if fname in title_dict:
        title_shape.text = "ファイル名" + fname + title_dict[fname]
    else:
        title_shape.text = "ファイル名" + fname
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)  # font size
    title_shape.text_frame.paragraphs[0].font.bold = True  # font bold
    
    tf = body_shape.text_frame
    tf.text = "頻出単語、品詞、出現回数"
    tf.paragraphs[0].font.size = Pt(18)  # font size
    tf.paragraphs[0].font.bold = True  # font bold
    for i in range(15):
        p = tf.add_paragraph()
        p.text = "「" + str(sorted_word_list[i]) + "(" + sorted_category_list[i] + ")」・・・" + str(sorted_num_list[i])
        p.level = 1
        p.font.size = Pt(14)  # font size
    
prs.save("result_filtered%s.pptx" % today)
